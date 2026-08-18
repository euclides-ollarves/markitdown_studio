"""
Test Word (.docx) and PowerPoint (.pptx) conversions
"""

import sys
import io
from pathlib import Path
import pptx
from pptx.util import Inches, Pt
import zipfile

root_dir = Path(__file__).resolve().parent.parent
if str(root_dir) not in sys.path:
    sys.path.insert(0, str(root_dir))

from gui.markitdown_service import MarkItDownService

def test_pptx_conversion():
    prs = pptx.Presentation()
    # Slide 1: Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "MarkItDown Presentation"
    subtitle.text = "Automated Document Pipeline Demo"

    # Slide 2: Bullet points
    bullet_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide2.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Key Benefits"
    tf = body_shape.text_frame
    tf.text = "Token efficient representation"
    p = tf.add_paragraph()
    p.text = "Native LLM comprehension"

    stream = io.BytesIO()
    prs.save(stream)
    pptx_bytes = stream.getvalue()

    service = MarkItDownService()
    out = service.convert_file(
        pptx_bytes,
        "presentation.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

    assert out.success, f"PPTX conversion failed: {out.error}"
    assert "MarkItDown Presentation" in out.markdown
    assert "Key Benefits" in out.markdown
    print(f" [PASS] PPTX Conversion -> {out.stats.word_count} words in {out.stats.duration_ms}ms")
    print(f" Preview:\n{out.markdown}\n")

if __name__ == "__main__":
    test_pptx_conversion()
